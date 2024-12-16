##### RAG Module #####

# Imports the sys module to access command-line arguments
import sys

# Filters warnings
import warnings

# Imports the listdir, isfile, join, and isdir functions from the os and os.path modules to handle directories and files
from os import listdir
from os.path import isdir, isfile, join

# Imports the docx module to handle Word files
import docx

# Imports the PyPDF2 module to handle PDF files
import PyPDF2

# Imports the HuggingFaceEmbeddings class from the langchain_huggingface package to create embeddings
from langchain_huggingface import HuggingFaceEmbeddings

# Imports the Qdrant class from the langchain_qdrant package to create a Qdrant instance and send data to the vector database
from langchain_qdrant import Qdrant

# Imports the TokenTextSplitter class from the langchain_text_splitters package to split text into tokens
from langchain_text_splitters import TokenTextSplitter

# Imports the Presentation module from the pptx package to handle PowerPoint files
from pptx import Presentation

# Imports the QdrantClient, Distance, and VectorParams classes from the qdrant_client package
# We'll create the Qdrant client and define storage parameters for the vector database
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams


import json


warnings.filterwarnings('ignore')

# Defines a function that lists all files in a directory, including those in subdirectories
def list_files(dir):

    # Initializes an empty list to store file paths
    file_list  = []

    # Iterates over all files and directories in the specified directory
    for f in listdir(dir):

        # If it's a file, adds it to the list
        if isfile(join(dir, f)):
            file_list .append(join(dir, f))

        # If it's a directory, recursively calls the function and adds the results to the list
        elif isdir(join(dir, f)):
            file_list  += list_files(join(dir, f))

    # Returns the list of files
    return file_list

# Defines a function that loads text from a Word file
def word_file_load(file_name):

    # Opens the Word file
    doc = docx.Document(file_name)

    # Extracts text from each paragraph and adds it to the list
    fullText = [para.text for para in doc.paragraphs]

    # Joins all text into a single string separated by line breaks
    return '\n'.join(fullText)

# Defines a function that loads text from a PowerPoint file
def pptx_file_load(file_name):

    # Opens the PowerPoint file
    prs = Presentation(file_name)

    # Initializes an empty list to store the text
    fullText = []

    # Iterates over all slides
    for slide in prs.slides:

        # Iterates over all shapes in the slide
        for shape in slide.shapes:

            # If the shape has a "text" attribute, adds the text to the list
            if hasattr(shape, "text"):
                fullText.append(shape.text)

    # Joins all text into a single string separated by line breaks
    return '\n'.join(fullText)


# Defines the main function for indexing documents
def main_indexing(mypath):

    # Defines the model name to be used for creating embeddings
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"

    # Defines the model configuration
    model_kwargs = {'device': 'cpu'}

    # Defines the encoding configuration
    encode_kwargs = {'normalize_embeddings': True}

    # Initializes the HuggingFace embeddings class
    hf = HuggingFaceEmbeddings(model_name=model_name,
                               model_kwargs=model_kwargs,
                               encode_kwargs=encode_kwargs)

    # Initializes the Qdrant client
    client = QdrantClient("http://localhost:6333")

    # Defines the name of the embeddings collection
    collection_name = "VectorDB"

    # If the collection already exists, deletes it
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    # Creates a new collection with specified parameters
    client.create_collection(collection_name,
                             vectors_config=VectorParams(size=768, distance=Distance.DOT))

    # Initializes the Qdrant instance
    qdrant = Qdrant(client, collection_name, hf)

    # Prints a message indicating that document indexing is starting
    print("\nIndexing documents...\n")

    # Gets the list of all files in the specified directory
    file_list = list_files(mypath)

    # Iterates over each file in the list
    for file in file_list:

        try:

            # Initializes an empty string to store the file content
            file_content = ""

            # Checks if the file is a PDF
            if file.endswith(".pdf"):

                print("Indexing: " + file)

                reader = PyPDF2.PdfReader(file)

                for page in reader.pages:
                    file_content += " " + page.extract_text()

            # Checks if the file is a plain text file
            elif file.endswith(".txt"):

                print("Indexing: " + file)

                with open(file, 'r') as f:
                    file_content = f.read()

            # Checks if the file is a Word file
            elif file.endswith(".docx"):

                print("Indexing: " + file)

                file_content = word_file_load(file)

            # Checks if the file is a PowerPoint file
            elif file.endswith(".pptx"):

                print("Indexing: " + file)

                file_content = pptx_file_load(file)

            else:

                # If the file format is unsupported, continues to the next file
                continue

            # Initializes the text splitter with specified chunk size and overlap
            text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)

            # Splits the file content into text chunks
            texts = text_splitter.split_text(file_content)

            # Creates metadata for each text chunk
            # This allows the LLM to reference the source
            metadata = [{"path": file} for _ in texts]

            # Adds the texts and their metadata to Qdrant
            qdrant.add_texts(texts, metadatas=metadata)

        except Exception as e:

            # If an error occurs, prints an error message
            print(f"Process failed for file {file}: {e}")

    # Prints a message indicating that indexing is complete
    print("\nIndexing Complete!\n")

# Checks if the script is being executed directly
if __name__ == "__main__":

    # Gets the command-line arguments
    arguments = sys.argv
    print (arguments)

    # Checks if a directory path was provided
    if len(arguments) > 1:
        main_indexing(arguments[1])
    else:
        # If not, prints an error message
        print("You need to provide a path to the folder with documents to index.")
